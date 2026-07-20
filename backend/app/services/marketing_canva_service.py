from __future__ import annotations

import io

import httpx
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.marketing_content_service import BRAND_LOGO_URL


async def _official_logo() -> bytes:
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.get(BRAND_LOGO_URL)
    response.raise_for_status()
    logo = Image.open(io.BytesIO(response.content)).convert("RGBA")
    bbox = logo.getbbox()
    if bbox is None:
        raise ValueError("Logo oficial sem conteudo visivel")
    output = io.BytesIO()
    logo.crop(bbox).save(output, format="PNG", optimize=True)
    return output.getvalue()


def _wrapped_title(title: str) -> str:
    words = " ".join(title.strip().split()).split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if not current or len(candidate) <= 21:
            current = candidate
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return "\n".join(lines[:4])


async def build_canva_pptx(art_content: bytes, title: str) -> bytes:
    """Gera um design 4:5 importavel no Canva com logo e titulo editaveis."""
    presentation = Presentation()
    presentation.slide_width = Inches(8)
    presentation.slide_height = Inches(10)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    # A imagem gerada permanece como fundo. A faixa esquerda encobre a versao
    # rasterizada e recebe os elementos nativos/editaveis do PPTX.
    slide.shapes.add_picture(io.BytesIO(art_content), 0, 0, width=Inches(8), height=Inches(10))
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.7), Inches(10))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(16, 0, 31)
    panel.line.fill.background()

    logo = await _official_logo()
    slide.shapes.add_picture(io.BytesIO(logo), Inches(0.52), Inches(0.52), width=Inches(2.25))

    accent = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.52),
        Inches(2.05),
        Inches(0.66),
        Inches(0.07),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(123, 0, 255)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.52), Inches(2.48), Inches(3.82), Inches(4.0))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = _wrapped_title(title)
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.name = "Inter"
    paragraph.font.size = Pt(34)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    paragraph.space_after = Pt(0)

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()
